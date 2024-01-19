import os
import PyPDF2
import logging
from pptx import Presentation

from model_server.config import logging_level


ROOT_DIRECTORY_PATH = os.getcwd()

logger = logging.getLogger(f"{__name__}")
logging.basicConfig()
logger.setLevel(logging_level)


async def pdf_extraction_alg(uploaded_file):

    unique_file_name = os.path.join(
        ROOT_DIRECTORY_PATH,
        uploaded_file.filename
    )
    logger.debug(f"Extracting information from PDF: {str(uploaded_file.filename)}")
    contents = await uploaded_file.read()

    with open(unique_file_name, 'wb') as destination_file:
        destination_file.write(contents)

    pdf_text = []

    with open(unique_file_name, 'rb') as pdf:
        reader = PyPDF2.PdfReader(pdf, strict=False)

        for page in reader.pages:
            content = page.extract_text()
            pdf_text.append(content)

    os.remove(unique_file_name)

    logger.debug(f"Success! Deleting: {str(uploaded_file.filename)}")

    output_string = ''

    for item in pdf_text:
        output_string += item

    return output_string


async def pptx_extraction_alg(uploaded_file):

    unique_file_name = os.path.join(
        ROOT_DIRECTORY_PATH,
        uploaded_file.filename
    )
    contents = await uploaded_file.read()

    with open(unique_file_name, 'wb') as destination_file:
        destination_file.write(contents)

    logger.debug(f"Extracting information from PPT file: {str(uploaded_file.filename)}")

    pres = Presentation(unique_file_name)

    output_string = ""

    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                temp = "\n\n" + shape.text
                output_string += temp

    return output_string
